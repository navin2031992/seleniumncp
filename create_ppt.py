from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime

# ── Colour palette ──────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW     = RGBColor(0xFF, 0xD1, 0x66)
RED        = RGBColor(0xFF, 0x4D, 0x6D)
GREEN      = RGBColor(0x52, 0xB7, 0x88)
ORANGE     = RGBColor(0xFF, 0x9F, 0x1C)
LIGHT_GREY = RGBColor(0xCC, 0xD6, 0xE0)
CARD_BG    = RGBColor(0x16, 0x2B, 0x3E)   # slightly lighter navy for cards

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

# ── Helpers ──────────────────────────────────────────────────
def blank_slide():
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)

def fill_bg(slide, color=DARK_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=CARD_BG, line_color=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h,
                 font_size=16, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def add_para(tf, text, font_size=14, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, space_before=Pt(4)):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    run = p.add_run()
    run.text = text
    run.font.size = _Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p

def accent_bar(slide, t=0.55, h=0.06):
    bar = slide.shapes.add_shape(1, 0, Inches(t), SLIDE_W, Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def divider(slide, t, color=ACCENT):
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(t), Inches(12.33), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()

def pill_label(slide, text, l, t, w=1.4, h=0.3, fill=ACCENT, font_color=DARK_BG, font_size=10):
    r = add_rect(slide, l, t, w, h, fill_color=fill)
    tf = r.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = font_color


# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)

accent_bar(s, t=0.4, h=0.08)

add_text_box(s, "SELENIUM MCP", 0.6, 0.7, 12, 1.0,
             font_size=52, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text_box(s, "+ JIRA INTEGRATION", 0.6, 1.65, 12, 0.7,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "AI-Driven Test Automation Framework", 0.6, 2.35, 12, 0.5,
             font_size=20, bold=False, color=ACCENT2, align=PP_ALIGN.CENTER, italic=True)

divider(s, t=3.0)

add_text_box(s, "Complete Analysis  |  Architecture  |  Risks  |  Automation Tester Role  |  Result Accuracy",
             0.6, 3.15, 12, 0.4, font_size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

add_text_box(s, f"Prepared by Navin  •  {datetime.date.today().strftime('%B %Y')}",
             0.6, 6.8, 12, 0.4, font_size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# corner decoration
add_rect(s, 0, 0, 0.15, 7.5, fill_color=ACCENT)
add_rect(s, 13.18, 0, 0.15, 7.5, fill_color=ACCENT)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "AGENDA", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

items = [
    ("01", "Project Overview & What is Selenium MCP"),
    ("02", "System Architecture & Agent Ecosystem"),
    ("03", "8 AI Agents Deep Dive"),
    ("04", "Selenium MCP Tools"),
    ("05", "Jira Integration — How It Works"),
    ("06", "End-to-End Pipeline Flow (7 Stages)"),
    ("07", "Where Human Automation Tester Is Required"),
    ("08", "Risks Involved"),
    ("09", "Will It Give 100% Results? — Honest Assessment"),
    ("10", "Benefits, Quality Gate & Recommendations"),
]

cols = [items[:5], items[5:]]
for ci, col in enumerate(cols):
    lx = 0.5 + ci * 6.5
    for ri, (num, title) in enumerate(col):
        ty = 1.45 + ri * 1.0
        add_rect(s, lx, ty, 0.55, 0.55, fill_color=ACCENT)
        add_text_box(s, num, lx, ty+0.02, 0.55, 0.5,
                     font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
        add_text_box(s, title, lx+0.65, ty+0.05, 5.5, 0.45,
                     font_size=13, color=WHITE)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — PROJECT OVERVIEW
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "PROJECT OVERVIEW", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

add_text_box(s,
    "Selenium MCP is a plug-and-play AI automation layer that sits on top of ANY existing test project.\n"
    "It uses Roo Code AI agents + Microsoft Edge + Jira to generate, run, fix and report tests — "
    "without replacing or imposing a new framework.",
    0.5, 1.4, 12.3, 0.9, font_size=14, color=ACCENT2)

cards = [
    (ACCENT,  "Language",    "Java • Python • JS/TS\nC# • Ruby — Any"),
    (GREEN,   "Browser",     "Microsoft Edge\n(Auto-managed driver)"),
    (YELLOW,  "AI Layer",    "Roo Code Agents\n(8 specialised modes)"),
    (ORANGE,  "Integration", "Jira Tickets →\nAuto Test Generation"),
    (RED,     "Protocol",    "MCP — Model\nContext Protocol"),
    (ACCENT2, "Node.js",     "v18+ Required\n@angiejones/mcp-selenium"),
]
for i, (col, title, body) in enumerate(cards):
    lx = 0.4 + i * 2.1
    add_rect(s, lx, 2.5, 2.0, 2.0, fill_color=CARD_BG, line_color=col, line_w=Pt(2))
    add_rect(s, lx, 2.5, 2.0, 0.35, fill_color=col)
    add_text_box(s, title, lx+0.05, 2.52, 1.9, 0.3,
                 font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text_box(s, body, lx+0.1, 2.9, 1.8, 1.4,
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 0.4, 5.0, 12.5, 1.6, fill_color=CARD_BG, line_color=ACCENT, line_w=Pt(1))
add_text_box(s, "KEY GOAL:", 0.6, 5.1, 2, 0.3, font_size=12, bold=True, color=YELLOW)
add_text_box(s,
    "Automate the entire testing lifecycle — from reading a Jira ticket to fixing broken tests — "
    "using AI agents with zero manual scripting, adapting to whatever test framework the team already uses.",
    0.6, 5.4, 12.1, 0.9, font_size=12, color=WHITE)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "SYSTEM ARCHITECTURE", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

# Three-layer diagram
layers = [
    (ACCENT,  "JIRA LAYER",     "Tickets  •  Acceptance Criteria  •  Sprint Backlog  •  Priority Mapping"),
    (YELLOW,  "AI AGENT LAYER", "8 Roo Code Agents  •  framework-profile.json  •  Pipeline Orchestrator"),
    (GREEN,   "SELENIUM LAYER", "MCP Server  •  Edge Browser  •  14 Selenium Tools  •  Screenshots"),
    (ORANGE,  "PROJECT LAYER",  "Your Existing Tests  •  Locator Registry  •  Fixtures  •  Reports"),
]
for i, (col, title, desc) in enumerate(layers):
    ty = 1.45 + i * 1.3
    add_rect(s, 0.4, ty, 12.5, 1.15, fill_color=CARD_BG, line_color=col, line_w=Pt(2))
    add_rect(s, 0.4, ty, 2.2, 1.15, fill_color=col)
    add_text_box(s, title, 0.4, ty+0.35, 2.2, 0.5,
                 font_size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text_box(s, desc, 2.8, ty+0.35, 9.8, 0.5,
                 font_size=12, color=WHITE)
    if i < 3:
        add_text_box(s, "▼", 6.3, ty+1.1, 0.5, 0.3,
                     font_size=14, color=ACCENT, align=PP_ALIGN.CENTER)

add_text_box(s,
    "Each layer is decoupled — change your test framework or Jira project without touching the others.",
    0.5, 6.85, 12.3, 0.4, font_size=11, color=ACCENT2, italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — 8 AGENTS OVERVIEW
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "8 AI AGENTS — THE CORE ENGINE", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

agents = [
    ("🔬", "Framework Analyzer",     "Scans project → writes framework-profile.json\nDetects language, framework, naming, structure",      ACCENT),
    ("🎫", "Jira Test Creator",       "Reads Jira ticket → generates tests in\nexisting project style with full AC coverage",              YELLOW),
    ("🔍", "XPath Discovery",         "Opens Edge → maps all UI elements →\nwrites locator registry JSON files",                           GREEN),
    ("▶️",  "Test Runner",            "Runs test suite → retries flaky tests →\ncaptures screenshots → writes HTML/JSON reports",          ORANGE),
    ("🔧", "Test Maintenance",        "Diagnoses failures → re-discovers broken\nselectors live → patches test files",                     RED),
    ("🔄", "Pipeline Orchestrator",   "Runs all 6 stages end-to-end →\nmanages state, recovery & exec summary",                           ACCENT2),
    ("📊", "Test Gap Analyzer",       "Cross-references Jira vs test files →\nfinds untested tickets sorted by priority",                  RGBColor(0xB5, 0x83, 0xFF)),
    ("📦", "Test Data Manager",       "Extracts hardcoded values → fixtures →\nflags security risks in test data",                        RGBColor(0xFF, 0x70, 0xA6)),
]

for i, (icon, name, desc, col) in enumerate(agents):
    col_i = i % 4
    row_i = i // 4
    lx = 0.3 + col_i * 3.25
    ty = 1.45 + row_i * 2.6
    add_rect(s, lx, ty, 3.1, 2.3, fill_color=CARD_BG, line_color=col, line_w=Pt(2))
    add_text_box(s, icon, lx+0.1, ty+0.08, 0.6, 0.5, font_size=22)
    add_text_box(s, name, lx+0.7, ty+0.1, 2.3, 0.4,
                 font_size=11, bold=True, color=col)
    add_text_box(s, desc, lx+0.1, ty+0.6, 2.9, 1.5,
                 font_size=10, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — SELENIUM MCP TOOLS
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "SELENIUM MCP — 14 BROWSER TOOLS", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

tools = [
    ("selenium_start_session",   "Open browser (Edge / Chrome / Firefox)"),
    ("selenium_navigate",        "Navigate to any URL"),
    ("selenium_find_element",    "Find element by id / css / xpath / name"),
    ("selenium_click",           "Click an element"),
    ("selenium_type",            "Type text into an input field"),
    ("selenium_get_text",        "Get visible text from element"),
    ("selenium_get_attribute",   "Read any HTML attribute value"),
    ("selenium_screenshot",      "Capture full-page screenshot (PNG)"),
    ("selenium_wait_for_element","Wait until element is visible/present"),
    ("selenium_execute_script",  "Run any JavaScript in browser context"),
    ("selenium_select_option",   "Choose from dropdown / select element"),
    ("selenium_hover",           "Hover over an element (tooltips, menus)"),
    ("selenium_get_page_source", "Get full HTML of current page"),
    ("selenium_close_session",   "Close browser and end session cleanly"),
]

col1 = tools[:7]
col2 = tools[7:]

for ci, col in enumerate([col1, col2]):
    lx = 0.4 + ci * 6.5
    for ri, (tool, desc) in enumerate(col):
        ty = 1.45 + ri * 0.73
        add_rect(s, lx, ty, 6.1, 0.65, fill_color=CARD_BG)
        pill_label(s, "TOOL", lx+0.08, ty+0.15, w=0.6, h=0.28, fill=ACCENT, font_size=8)
        add_text_box(s, tool, lx+0.76, ty+0.05, 2.5, 0.3,
                     font_size=10, bold=True, color=ACCENT2)
        add_text_box(s, desc, lx+0.76, ty+0.35, 5.1, 0.28,
                     font_size=9.5, color=LIGHT_GREY)

add_text_box(s,
    "Package: @angiejones/mcp-selenium v0.2.3  •  Runs via npx  •  Auto-starts with Roo Code",
    0.4, 6.9, 12.5, 0.35, font_size=11, color=ACCENT2, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 7 — JIRA INTEGRATION
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "JIRA INTEGRATION — HOW IT WORKS", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

add_text_box(s,
    "Jira is accessed via the corporate LLM's built-in integration — NO separate Jira MCP package.\n"
    "Agents read tickets directly and map Jira issue types and priorities to test categories.",
    0.5, 1.4, 12.3, 0.7, font_size=13, color=ACCENT2)

# Left — What gets read
add_rect(s, 0.4, 2.2, 5.8, 3.8, fill_color=CARD_BG, line_color=YELLOW, line_w=Pt(1.5))
add_text_box(s, "WHAT JIRA PROVIDES", 0.5, 2.3, 5.6, 0.35,
             font_size=13, bold=True, color=YELLOW)
jira_items = [
    "• Ticket Summary → test suite / describe block name",
    "• Acceptance Criteria → each criterion = ≥1 test case",
    "• Issue Type → Story=E2E, Bug=Regression, Task=Functional",
    "• Priority → Blocker/Critical=P0, High=P1, Medium=P2",
    "• Sprint → batch generate tests per sprint",
    "• Labels/Components → which pages to test",
    "• @jira TICKET-ID tag → full traceability in test files",
]
for i, item in enumerate(jira_items):
    add_text_box(s, item, 0.6, 2.7+i*0.42, 5.4, 0.38, font_size=11, color=WHITE)

# Right — Priority mapping
add_rect(s, 6.5, 2.2, 6.4, 1.7, fill_color=CARD_BG, line_color=ORANGE, line_w=Pt(1.5))
add_text_box(s, "PRIORITY MAPPING", 6.6, 2.3, 6.2, 0.35,
             font_size=13, bold=True, color=ORANGE)
pmap = [("Blocker / Critical", "→ P0", RED), ("High", "→ P1", ORANGE),
        ("Medium", "→ P2", YELLOW), ("Low", "→ P3", GREEN)]
for i, (src, tgt, col) in enumerate(pmap):
    ty = 2.7 + i * 0.33
    add_text_box(s, src, 6.6, ty, 3.2, 0.3, font_size=11, color=LIGHT_GREY)
    add_text_box(s, tgt, 9.8, ty, 2.8, 0.3, font_size=11, bold=True, color=col)

# Right — Issue type mapping
add_rect(s, 6.5, 4.05, 6.4, 2.0, fill_color=CARD_BG, line_color=GREEN, line_w=Pt(1.5))
add_text_box(s, "ISSUE TYPE → TEST TYPE", 6.6, 4.15, 6.2, 0.35,
             font_size=13, bold=True, color=GREEN)
imap = [("Story", "→ E2E Tests"), ("Bug", "→ Regression Tests"),
        ("Task", "→ Functional Tests"), ("Sub-task", "→ Unit Tests"), ("Epic", "→ Test Suite")]
for i, (src, tgt) in enumerate(imap):
    ty = 4.55 + i * 0.28
    add_text_box(s, src, 6.6, ty, 3.0, 0.26, font_size=10.5, color=LIGHT_GREY)
    add_text_box(s, tgt, 9.6, ty, 3.0, 0.26, font_size=10.5, bold=True, color=ACCENT2)

add_text_box(s,
    "Config: config/jira.config.json  •  Credentials: env vars JIRA_BASE_URL / JIRA_EMAIL / JIRA_API_TOKEN",
    0.4, 6.4, 12.5, 0.35, font_size=10, color=LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — PIPELINE FLOW
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "END-TO-END PIPELINE — 7 STAGES", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

stages = [
    ("0", "Framework\nAnalysis",   "Scan project\nWrite profile",            ACCENT,  "ONCE"),
    ("1", "Jira\nIntake",          "Fetch tickets\nACs & priority",           YELLOW,  "PER SPRINT"),
    ("2", "Locator\nDiscovery",    "Open Edge\nMap elements",                 GREEN,   "PER PAGE"),
    ("3", "Test\nGeneration",      "AI writes tests\nin project style",       ORANGE,  "PER TICKET"),
    ("4", "Execution",             "Run suite\nCapture failures",             RED,     "EVERY RUN"),
    ("5", "Maintenance",           "Heal selectors\nFix broken tests",        ACCENT2, "ON FAILURE"),
    ("6", "Reporting",             "Executive summary\nQuality gate",         RGBColor(0xB5,0x83,0xFF), "EVERY RUN"),
]

for i, (num, title, desc, col, freq) in enumerate(stages):
    lx = 0.3 + i * 1.87
    add_rect(s, lx, 1.5, 1.75, 3.5, fill_color=CARD_BG, line_color=col, line_w=Pt(2))
    add_rect(s, lx, 1.5, 1.75, 0.5, fill_color=col)
    add_text_box(s, f"STAGE {num}", lx+0.05, 1.53, 1.65, 0.42,
                 font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text_box(s, title, lx+0.1, 2.08, 1.55, 0.7,
                 font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text_box(s, desc, lx+0.1, 2.85, 1.55, 1.0,
                 font_size=10, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    pill_label(s, freq, lx+0.17, 4.75, w=1.4, h=0.25, fill=col, font_color=DARK_BG, font_size=8)
    if i < 6:
        add_text_box(s, "→", lx+1.75, 2.9, 0.12, 0.4,
                     font_size=18, color=ACCENT, align=PP_ALIGN.CENTER)

add_rect(s, 0.3, 5.2, 12.7, 1.0, fill_color=CARD_BG, line_color=ACCENT, line_w=Pt(1))
add_text_box(s, "STATE FILE:", 0.5, 5.3, 1.5, 0.3, font_size=11, bold=True, color=YELLOW)
add_text_box(s,
    "tests/pipeline-state.json tracks every stage status (pending → in_progress → completed → failed).\n"
    "Pipeline can resume from any failed stage. Non-blocking failures are logged and pipeline continues.",
    2.0, 5.3, 10.8, 0.7, font_size=11, color=WHITE)


# ════════════════════════════════════════════════════════════
# SLIDE 9 — WHERE HUMAN TESTER IS REQUIRED
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "WHERE AUTOMATION TESTER IS REQUIRED", 0.5, 0.7, 12, 0.5,
             font_size=26, bold=True, color=ACCENT)
divider(s, 1.3)

add_text_box(s,
    "This system is AI-assisted, NOT fully autonomous. Skilled automation testers are essential at multiple points.",
    0.5, 1.4, 12.3, 0.4, font_size=13, color=ACCENT2, italic=True)

human_roles = [
    (YELLOW, "SETUP & CONFIGURATION",
     "• Install Node.js 18+, configure .env file\n"
     "• Set JIRA credentials (JIRA_BASE_URL, JIRA_EMAIL, JIRA_API_TOKEN)\n"
     "• Configure environments.json with real application URLs\n"
     "• Configure selenium.config.json (timeouts, retry count, quality gates)"),

    (ORANGE, "FRAMEWORK ANALYSIS REVIEW",
     "• Review framework-profile.json after Framework Analyzer runs\n"
     "• Correct any misidentified framework or wrong run command\n"
     "• Add missing code templates the AI may have missed\n"
     "• Verify folder structure mappings are accurate"),

    (RED, "TEST QUALITY REVIEW",
     "• Review generated tests for logical correctness (not just syntax)\n"
     "• Validate that acceptance criteria are covered correctly\n"
     "• Check boundary values and negative test cases are meaningful\n"
     "• Ensure Page Object gaps are implemented before execution"),

    (GREEN, "MAINTENANCE ESCALATION",
     "• Triage failures the AI marks as 'manual review required'\n"
     "• Decide: real bug vs test issue vs environment problem\n"
     "• Implement fixes for complex selector or logic changes\n"
     "• Approve healed tests before re-running in CI/CD"),

    (ACCENT2, "CI/CD & INFRASTRUCTURE",
     "• Set up headless mode (SELENIUM_HEADLESS=true) for CI pipelines\n"
     "• Ensure correct msedgedriver version on CI machines\n"
     "• Configure webhook URLs for pipeline notifications\n"
     "• Monitor flakiness metrics and adjust retry thresholds"),

    (RGBColor(0xB5,0x83,0xFF), "SECURITY & DATA AUDIT",
     "• Review Test Data Manager security audit reports\n"
     "• Approve or reject flagged HIGH/MEDIUM risk credentials\n"
     "• Ensure no real credentials leak into fixture files\n"
     "• Sign off on the test data strategy per environment"),
]

for i, (col, title, body) in enumerate(human_roles):
    ci = i % 3
    ri = i // 3
    lx = 0.3 + ci * 4.35
    ty = 1.9 + ri * 2.35
    add_rect(s, lx, ty, 4.2, 2.15, fill_color=CARD_BG, line_color=col, line_w=Pt(2))
    add_rect(s, lx, ty, 4.2, 0.32, fill_color=col)
    add_text_box(s, title, lx+0.08, ty+0.04, 4.0, 0.25,
                 font_size=9.5, bold=True, color=DARK_BG)
    add_text_box(s, body, lx+0.1, ty+0.4, 3.95, 1.6,
                 font_size=9, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — RISKS INVOLVED
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "RISKS INVOLVED", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=RED)
divider(s, 1.3)

risks = [
    ("HIGH", RED,
     "AI-Generated Tests May Miss Edge Cases",
     "AI generates tests from ACs only. Domain-specific edge cases, regulatory rules, or implicit business logic NOT in the ticket will be missed. A human tester must review every generated test file."),

    ("HIGH", RED,
     "Selector Fragility After UI Changes",
     "Locators discovered today may break after any UI redesign. XPath Discovery must be re-run after every deployment that touches the frontend. Relying on position-based XPath is especially risky."),

    ("HIGH", RED,
     "Jira Integration Dependency",
     "The system relies on corporate LLM's built-in Jira integration. If Jira is unavailable or the LLM cannot access it, test generation stops. Manual ticket pasting is the fallback — adds friction."),

    ("MEDIUM", ORANGE,
     "Driver Version Mismatch",
     "Edge browser auto-updates silently. If msedgedriver version lags behind (e.g., old v133 vs new v148), all sessions fail instantly. Requires selenium-manager to be active and cache to be writable."),

    ("MEDIUM", ORANGE,
     "Flaky Test Rate",
     "The quality gate flags >5% flakiness. Network timeouts, slow CI machines, or race conditions in SPAs can inflate flakiness. Retry logic (3×) helps but adds runtime. Root cause must be fixed."),

    ("MEDIUM", ORANGE,
     "Framework Profile Accuracy",
     "If framework-profile.json is wrong (wrong run command, wrong folder), ALL generated tests will be in the wrong location or fail to execute. Profile must be manually verified after each project scan."),

    ("LOW", YELLOW,
     "Security — Hardcoded Credentials",
     "Test Data Manager flags HIGH risk credentials. If an engineer ignores the audit, real passwords or API keys could be committed to git inside fixture files. Policy enforcement is a human responsibility."),

    ("LOW", YELLOW,
     "100% Pass Rate Is Not Guaranteed",
     "Quality gate is set to 90% minimum pass rate. Newly generated tests will have failures until ACs are refined, Page Objects are implemented, and the test environment is stable."),
]

for i, (level, col, title, body) in enumerate(risks):
    ci = i % 2
    ri = i // 2
    lx = 0.3 + ci * 6.5
    ty = 1.4 + ri * 1.45
    add_rect(s, lx, ty, 6.3, 1.3, fill_color=CARD_BG, line_color=col, line_w=Pt(1.5))
    pill_label(s, level, lx+0.12, ty+0.12, w=0.75, h=0.26, fill=col, font_color=DARK_BG, font_size=9)
    add_text_box(s, title, lx+1.0, ty+0.1, 5.1, 0.3,
                 font_size=11, bold=True, color=col)
    add_text_box(s, body, lx+0.12, ty+0.5, 6.0, 0.72,
                 font_size=9.5, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════
# SLIDE 11 — 100% RESULT ASSESSMENT
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "WILL IT GIVE 100% RESULTS?", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=YELLOW)
divider(s, 1.3)

add_text_box(s, "HONEST ASSESSMENT", 0.5, 1.4, 12, 0.35,
             font_size=14, bold=True, color=ACCENT2)

# Big verdict
add_rect(s, 0.4, 1.85, 12.5, 0.95, fill_color=RGBColor(0x3A, 0x10, 0x10), line_color=RED, line_w=Pt(2))
add_text_box(s,
    "NO — 100% automated results without human review is NOT achievable with this or any AI test automation system.",
    0.6, 1.95, 12.1, 0.65, font_size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)

cols_data = [
    ("WHAT IT DOES WELL  ✅", GREEN, [
        "Generates syntactically correct test code rapidly (minutes vs days)",
        "Covers happy path + primary negative cases from ACs consistently",
        "Auto-heals broken selectors by re-discovering them on live pages",
        "Batch-processes entire sprints without manual scripting",
        "Produces structured reports with pass rate and quality gate",
        "Eliminates repetitive boilerplate test writing",
        "Maintains locator registry reducing selector maintenance time by ~70%",
    ]),
    ("WHAT STILL NEEDS HUMANS  ⚠️", ORANGE, [
        "Domain logic & implicit business rules not in Jira tickets",
        "Complex multi-step scenarios spanning multiple tickets",
        "Performance, load, security and accessibility testing",
        "Page Object methods must be manually implemented after generation",
        "Assertion correctness — AI may assert wrong expected values",
        "Test data strategy for environment-specific values",
        "Final sign-off before any test is run in production",
    ]),
]

for ci, (heading, col, items) in enumerate(cols_data):
    lx = 0.3 + ci * 6.5
    add_rect(s, lx, 2.95, 6.4, 3.9, fill_color=CARD_BG, line_color=col, line_w=Pt(1.5))
    add_text_box(s, heading, lx+0.15, 3.05, 6.1, 0.35,
                 font_size=12, bold=True, color=col)
    for ri, item in enumerate(items):
        add_text_box(s, f"• {item}", lx+0.2, 3.48+ri*0.47, 6.0, 0.42,
                     font_size=10, color=WHITE)

add_rect(s, 0.4, 7.05, 12.5, 0.38, fill_color=CARD_BG, line_color=YELLOW, line_w=Pt(1))
add_text_box(s,
    "REALISTIC EXPECTATION: 80–90% pass rate at first run per sprint. Rises to 95%+ after 2–3 iterations of human review + healing.",
    0.6, 7.1, 12.1, 0.3, font_size=11, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 12 — BENEFITS & QUALITY GATE
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "BENEFITS & QUALITY GATE", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

benefits = [
    (ACCENT,  "Speed",         "Test generation in minutes\ninstead of days per sprint"),
    (GREEN,   "Coverage",      "Every AC → at least one test\nBug tickets get regression tests"),
    (YELLOW,  "Traceability",  "@jira TICKET-ID in every test\nEnd-to-end audit trail"),
    (ORANGE,  "Self-Healing",  "Broken selectors auto-fixed\nNo manual locator maintenance"),
    (RED,     "Reporting",     "HTML + JSON + JUnit reports\nCI/CD ready exit codes"),
    (ACCENT2, "Framework",     "Works with ANY existing project\nNo framework migration"),
]

for i, (col, title, desc) in enumerate(benefits):
    ci = i % 3
    ri = i // 3
    lx = 0.3 + ci * 4.35
    ty = 1.45 + ri * 1.55
    add_rect(s, lx, ty, 4.2, 1.35, fill_color=CARD_BG, line_color=col, line_w=Pt(2))
    add_rect(s, lx, ty, 0.18, 1.35, fill_color=col)
    add_text_box(s, title, lx+0.3, ty+0.08, 3.7, 0.35,
                 font_size=13, bold=True, color=col)
    add_text_box(s, desc, lx+0.3, ty+0.5, 3.7, 0.7,
                 font_size=11, color=LIGHT_GREY)

# Quality gate box
add_rect(s, 0.4, 4.65, 12.5, 2.0, fill_color=CARD_BG, line_color=ACCENT, line_w=Pt(1.5))
add_text_box(s, "QUALITY GATE CONFIGURATION (selenium.config.json)",
             0.6, 4.73, 12, 0.35, font_size=13, bold=True, color=ACCENT)

qg = [
    ("minPassRate: 90%",   "Pipeline WARNS if pass rate < 90%. Release should be blocked."),
    ("maxFlakiness: 5%",   "More than 5% flaky tests triggers a maintenance review."),
    ("retries: 3×",        "Each failing test is retried with 2s/4s/8s exponential backoff."),
    ("Exit code 0/1/2/3",  "CI/CD gets 0=pass, 1=fail, 2=pipeline error, 3=env unreachable."),
]
for i, (setting, desc) in enumerate(qg):
    lx = 0.6 + (i % 2) * 6.2
    ty = 5.12 + (i // 2) * 0.42
    add_text_box(s, f"• {setting}", lx, ty, 2.2, 0.38,
                 font_size=11, bold=True, color=YELLOW)
    add_text_box(s, desc, lx+2.2, ty, 3.7, 0.38,
                 font_size=11, color=WHITE)


# ════════════════════════════════════════════════════════════
# SLIDE 13 — SETUP & QUICK START
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "SETUP & QUICK START", 0.5, 0.7, 12, 0.5,
             font_size=28, bold=True, color=ACCENT)
divider(s, 1.3)

steps = [
    (ACCENT,  "Step 1", "PREREQUISITES",
     "Node.js 18+  •  Microsoft Edge 148+  •  Roo Code VS Code extension\nnpm install in c:\\NewInitiatives\\seleniummcp"),
    (YELLOW,  "Step 2", "CONFIGURE",
     "Copy .env.example → .env  •  Set TEST_ENV, SELENIUM_BROWSER=edge\nSet JIRA_BASE_URL, JIRA_EMAIL, JIRA_API_TOKEN in .env"),
    (GREEN,   "Step 3", "VERIFY",
     "node verify-setup.js  →  14 checks must all pass (✅)\nOpen Roo Code → ask any agent 'What MCP tools are available?'"),
    (ORANGE,  "Step 4", "ANALYZE",
     "Switch to 🔬 Framework Analyzer\nSay: 'Scan my test project at {path}' → .roo/framework-profile.json created"),
    (RED,     "Step 5", "DISCOVER",
     "Switch to 🔍 XPath Discovery\nSay: 'Discover locators for https://your-app.com/login'  →  Edge opens"),
    (ACCENT2, "Step 6", "GENERATE",
     "Switch to 🎫 Jira Test Creator\nSay: 'Generate tests for PROJ-123'  →  test file created in project"),
    (YELLOW,  "Step 7", "EXECUTE",
     "Switch to ▶️ Test Runner\nSay: 'Run smoke tests on staging'  →  report at tests/reports/"),
    (RGBColor(0xB5,0x83,0xFF), "OR", "FULL AUTO",
     "Switch to 🔄 Pipeline Orchestrator\nSay: 'Run full pipeline for Sprint 42'  →  all stages automated"),
]

for i, (col, label, title, body) in enumerate(steps):
    ci = i % 2
    ri = i // 2
    lx = 0.3 + ci * 6.5
    ty = 1.45 + ri * 1.45
    add_rect(s, lx, ty, 6.3, 1.3, fill_color=CARD_BG, line_color=col, line_w=Pt(1.5))
    add_rect(s, lx, ty, 0.85, 1.3, fill_color=col)
    add_text_box(s, label, lx+0.02, ty+0.2, 0.8, 0.5,
                 font_size=11, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_text_box(s, title, lx+1.0, ty+0.06, 5.1, 0.3,
                 font_size=11, bold=True, color=col)
    add_text_box(s, body, lx+1.0, ty+0.43, 5.2, 0.78,
                 font_size=10, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════
# SLIDE 14 — RECOMMENDATIONS & CONCLUSION
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)
accent_bar(s)

add_text_box(s, "RECOMMENDATIONS & CONCLUSION", 0.5, 0.7, 12, 0.5,
             font_size=26, bold=True, color=ACCENT)
divider(s, 1.3)

add_rect(s, 0.4, 1.4, 12.5, 2.3, fill_color=CARD_BG, line_color=GREEN, line_w=Pt(1.5))
add_text_box(s, "RECOMMENDATIONS FOR MAXIMUM EFFECTIVENESS", 0.6, 1.5, 12, 0.35,
             font_size=13, bold=True, color=GREEN)
recs = [
    "1. Always run Framework Analyzer first on any new project — quality of profile determines quality of all generated tests.",
    "2. Assign a senior automation tester to review generated tests each sprint — AI accuracy is ~80% without review.",
    "3. Re-run XPath Discovery after every UI release — stale locators are the #1 cause of test failures.",
    "4. Set SELENIUM_HEADLESS=true in CI/CD — headed mode is for local debugging only.",
    "5. Monitor health-metrics.json weekly — rising flakiness signals environment or selector drift before it becomes critical.",
    "6. Treat Test Data Manager security audit as mandatory — never merge test code with HIGH-risk credential findings.",
]
for i, rec in enumerate(recs):
    ty = 1.9 + i * 0.3
    add_text_box(s, rec, 0.6, ty, 12.1, 0.27, font_size=10.5, color=WHITE)

add_rect(s, 0.4, 3.85, 12.5, 2.6, fill_color=CARD_BG, line_color=ACCENT, line_w=Pt(1.5))
add_text_box(s, "CONCLUSION", 0.6, 3.95, 12, 0.35,
             font_size=13, bold=True, color=ACCENT)

conclusion_paras = [
    ("This Selenium MCP + Jira integration is a state-of-the-art AI automation framework that dramatically "
     "accelerates the testing lifecycle. It is not a replacement for skilled automation testers — it is a "
     "force multiplier that removes the most tedious parts of their job.", WHITE, False),
    ("With proper setup, framework profile verification, and human review at key checkpoints, teams can realistically "
     "achieve 90–95% test pass rates per sprint within 2–3 cycles of usage.", ACCENT2, False),
    ("The 8 AI agents, 14 Selenium MCP tools, and 7-stage pipeline create a complete, "
     "language-agnostic automation ecosystem that grows smarter with every sprint.", LIGHT_GREY, True),
]
for i, (text, col, italic) in enumerate(conclusion_paras):
    ty = 4.38 + i * 0.6
    add_text_box(s, text, 0.6, ty, 12.1, 0.55, font_size=11, color=col, italic=italic)

add_text_box(s,
    "Project: c:\\NewInitiatives\\mcp\\seleniummcp  •  Verified: 2026-06-03  •  All 14 checks PASSED",
    0.4, 6.85, 12.5, 0.35, font_size=10, color=LIGHT_GREY, italic=True, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 15 — THANK YOU
# ════════════════════════════════════════════════════════════
s = blank_slide()
fill_bg(s)

add_rect(s, 0, 0, 0.15, 7.5, fill_color=ACCENT)
add_rect(s, 13.18, 0, 0.15, 7.5, fill_color=ACCENT)
accent_bar(s, t=3.3, h=0.06)

add_text_box(s, "THANK YOU", 0.6, 1.5, 12, 1.2,
             font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "Selenium MCP + Jira Integration", 0.6, 2.8, 12, 0.6,
             font_size=22, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

divider(s, t=3.5)

add_text_box(s,
    "Questions & Discussion",
    0.6, 3.7, 12, 0.5, font_size=20, color=ACCENT2, align=PP_ALIGN.CENTER)

details = [
    "Project Path:   c:\\NewInitiatives\\mcp\\seleniummcp",
    "Browser:        Microsoft Edge 148  (auto-managed driver)",
    "Package:        @angiejones/mcp-selenium v0.2.3",
    "Agents:         8 Roo Code AI Agents",
    "Last Verified:  2026-06-03 — ALL 14 CHECKS PASSED",
]
for i, d in enumerate(details):
    add_text_box(s, d, 3.0, 4.4+i*0.37, 7.5, 0.33,
                 font_size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────
out = r"c:\NewInitiatives\mcp\seleniummcp\Selenium_MCP_Jira_Integration.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
